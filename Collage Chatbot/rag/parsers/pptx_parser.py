import io
import hashlib
from typing import Dict, Any, Optional, List

class PPTXParser:
    """
    PPTX document parser for extracting text, metadata, and structure from PowerPoint presentations.
    Supports basic text extraction, slide-level processing, and metadata extraction.
    """
    
    def __init__(self):
        self.supported_formats = ['pptx']
        self.pptx_available = False
        
        try:
            from pptx import Presentation
            self.Presentation = Presentation
            self.pptx_available = True
        except ImportError:
            print("[PPTXParser] python-pptx not installed. Install with: pip install python-pptx")
    
    def parse_pptx(self, file_bytes: bytes, filename: str = "") -> Dict[str, Any]:
        """
        Parse PPTX file and extract text, metadata, and structure.
        
        Args:
            file_bytes: Raw bytes of the PPTX file
            filename: Original filename of the PPTX
            
        Returns:
            Dictionary containing extracted content and metadata
        """
        if not self.pptx_available:
            return {
                "success": False,
                "error": "python-pptx not installed",
                "filename": filename
            }
        
        try:
            ppt_file = io.BytesIO(file_bytes)
            prs = self.Presentation(ppt_file)
            
            # Extract basic metadata
            metadata = self._extract_metadata(prs, filename)
            
            # Extract text from all slides
            slides_content = []
            full_text = ""
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = ""
                slide_shapes = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        shape_text = shape.text.strip()
                        slide_text += shape_text + "\n"
                        slide_shapes.append({
                            "shape_type": str(shape.shape_type),
                            "text": shape_text,
                            "char_count": len(shape_text)
                        })
                
                if slide_text.strip():
                    slide_content = {
                        "slide_number": slide_num + 1,
                        "text": slide_text.strip(),
                        "shapes": slide_shapes,
                        "char_count": len(slide_text)
                    }
                    slides_content.append(slide_content)
                    full_text += slide_text + "\n\n"
            
            # Clean and structure the extracted text
            clean_text = self._clean_text(full_text)
            
            # Generate content hash
            content_hash = hashlib.sha256(clean_text.encode('utf-8')).hexdigest()
            
            # Extract sections (each slide is a section)
            sections = self._extract_sections(slides_content)
            
            return {
                "success": True,
                "filename": filename,
                "format": "pptx",
                "metadata": metadata,
                "full_text": clean_text,
                "content_hash": content_hash,
                "slides": slides_content,
                "total_slides": len(slides_content),
                "total_characters": len(clean_text),
                "sections": sections,
                "word_count": len(clean_text.split())
            }
            
        except Exception as e:
            print(f"[PPTXParser] Error parsing PPTX: {e}")
            return {
                "success": False,
                "error": str(e),
                "filename": filename
            }
    
    def _extract_metadata(self, prs, filename: str) -> Dict[str, Any]:
        """Extract metadata from PPTX presentation"""
        metadata = {
            "filename": filename,
            "title": "",
            "author": "",
            "subject": "",
            "creator": "",
            "created": "",
            "modified": "",
            "last_modified_by": "",
            "revision": "",
            "slide_count": len(prs.slides)
        }
        
        try:
            core_props = prs.core_properties
            if core_props:
                metadata.update({
                    "title": core_props.title or "",
                    "author": core_props.author or "",
                    "subject": core_props.subject or "",
                    "creator": core_props.creator or "",
                    "created": str(core_props.created) if core_props.created else "",
                    "modified": str(core_props.modified) if core_props.modified else "",
                    "last_modified_by": core_props.last_modified_by or "",
                    "revision": str(core_props.revision) if core_props.revision else ""
                })
        except Exception as e:
            print(f"[PPTXParser] Error extracting metadata: {e}")
        
        return metadata
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text"""
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = ' '.join(text.split())
        
        # Normalize line breaks
        text = text.replace('\r\n', '\n').replace('\r', '\n')
        
        # Remove multiple consecutive newlines
        while '\n\n\n' in text:
            text = text.replace('\n\n\n', '\n\n')
        
        return text.strip()
    
    def _extract_sections(self, slides_content: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """
        Extract sections from slides. Each slide is treated as a section.
        """
        sections = []
        
        for slide in slides_content:
            # Use first 50 chars as title if available
            title = slide["text"][:50] + "..." if len(slide["text"]) > 50 else slide["text"]
            if not title:
                title = f"Slide {slide['slide_number']}"
            
            sections.append({
                "title": title,
                "content": slide["text"],
                "slide_number": slide["slide_number"]
            })
        
        return sections
    
    def is_supported_format(self, filename: str) -> bool:
        """Check if the file format is supported"""
        return filename.lower().endswith('.pptx')
    
    def validate_pptx(self, file_bytes: bytes):
        """
        Validate if the file is a valid PPTX.
        
        Returns:
            Tuple of (is_valid, error_message)
        """
        if not self.pptx_available:
            return False, "python-pptx not installed"
        
        try:
            ppt_file = io.BytesIO(file_bytes)
            prs = self.Presentation(ppt_file)
            
            # Try to read the first slide to validate
            if len(prs.slides) > 0:
                first_slide = prs.slides[0]
                _ = [shape.text for shape in first_slide.shapes if hasattr(shape, "text")]
                
            return True, ""
            
        except Exception as e:
            return False, f"Invalid PPTX file: {str(e)}"
